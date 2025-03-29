from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import requests
import wikipedia
import spacy
from gtts import gTTS
from fastapi.middleware.cors import CORSMiddleware
from io import BytesIO
import os

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

nlp = spacy.load("en_core_web_sm")

PEXELS_API_KEY = "mYGJaFjke2BYgo5GtFZHP5HbAS5TGj5ad2LX5EPbIPeb7FaMZIkTRreK"

class TopicInput(BaseModel):
    topic: str
    subtopics: list[str] = []
    bg_color: str = "#FFFFFF"
    text_color: str = "#000000"
    font_style: str = "Times New Roman"

def fetch_wikipedia_content(topic):
    try:
        summary = wikipedia.summary(topic, sentences=5)
        content = summary.split(". ")
        return [point.strip() for point in content if len(point) > 20][:4]
    except Exception:
        return ["No relevant data found."]

def fetch_images(query, num_images=1):
    headers = {"Authorization": PEXELS_API_KEY}
    params = {"query": query, "per_page": num_images}
    response = requests.get("https://api.pexels.com/v1/search", headers=headers, params=params)
    if response.status_code == 200:
        return [photo["src"]["medium"] for photo in response.json().get("photos", [])]
    return []

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def generate_audio(text, filename):
    tts = gTTS(text)
    audio_folder = "generated_audios"
    os.makedirs(audio_folder, exist_ok=True)
    audio_path = os.path.join(audio_folder, filename)
    tts.save(audio_path)
    return audio_path

@app.post("/generate_ppt/")
def generate_ppt(input_data: TopicInput):
    topic = input_data.topic.strip()
    if not topic:
        raise HTTPException(status_code=400, detail="Topic is required")

    presentation = Presentation()
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)

    bg_rgb = hex_to_rgb(input_data.bg_color)
    text_rgb = hex_to_rgb(input_data.text_color)

    subtopics = input_data.subtopics or ["Introduction", "Benefits", "Challenges", "Future Scope"]

    topic_images = fetch_images(topic, num_images=len(subtopics))

    for idx, subtopic in enumerate(subtopics):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_rgb

        title = slide.shapes.title
        title.text = f"{subtopic} - {topic}"
        title.text_frame.paragraphs[0].font.color.rgb = text_rgb
        title.text_frame.paragraphs[0].font.name = input_data.font_style

        text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(4))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        content_points = fetch_wikipedia_content(f"{topic} {subtopic}")
        content_text = "\n".join(content_points)

        text_frame.text = content_text
        text_frame.paragraphs[0].font.color.rgb = text_rgb
        text_frame.paragraphs[0].font.name = input_data.font_style

        audio_file = generate_audio(content_text, f"{topic.replace(' ', '_')}_slide{idx + 1}.mp3")
        slide.shapes.add_movie(audio_file, Inches(1), Inches(6), Inches(2), Inches(2))

        if idx < len(topic_images):
            try:
                img_data = requests.get(topic_images[idx]).content
                image_stream = BytesIO(img_data)
                slide.shapes.add_picture(image_stream, Inches(11), Inches(2), width=Inches(4))
            except Exception as e:
                print(f"Error loading image {idx}: {e}")

    ppt_folder = "generated_ppts"
    os.makedirs(ppt_folder, exist_ok=True)
    ppt_path = os.path.join(ppt_folder, f"{topic.replace(' ', '_')}.pptx")
    presentation.save(ppt_path)

    return {"ppt_url": f"http://localhost:8000/download/{topic.replace(' ', '_')}.pptx"}

@app.get("/download/{filename}")
async def download_ppt(filename: str):
    ppt_folder = "generated_ppts"
    file_path = os.path.join(ppt_folder, filename)
    if os.path.exists(file_path):
        return FileResponse(file_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=filename)
    raise HTTPException(status_code=404, detail="File not found")